import streamlit as st
import google.generativeai as genai
from groq import Groq
from PIL import Image
import io, base64, time, requests, sys
import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from google.api_core import exceptions

# --- ضبط الترميز لدعم اللغة العربية بشكل كامل ---
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')

# --- 1. إعدادات الصفحة والتصميم العصري ---
st.set_page_config(page_title="AI Architect Ultimate", page_icon="🚀", layout="wide")

st.markdown("""
    <style>
    .stApp { background: linear-gradient(135deg, #0f172a 0%, #1e1e2f 100%); color: #ffffff; }
    .stTabs [data-baseweb="tab-list"] { gap: 10px; background-color: transparent; }
    .stTabs [data-baseweb="tab"] { height: 50px; background-color: rgba(255, 255, 255, 0.05); border-radius: 10px; color: white; font-weight: bold; }
    .stButton>button { background: linear-gradient(90deg, #00d2ff 0%, #3a7bd5 100%); color: white; border: none; padding: 12px; border-radius: 12px; font-weight: 700; width: 100%; transition: 0.3s; }
    .stButton>button:hover { transform: translateY(-2px); box-shadow: 0px 10px 20px rgba(0, 210, 255, 0.3); }
    .result-box { background: rgba(255, 255, 255, 0.03); backdrop-filter: blur(10px); padding: 20px; border-radius: 15px; border-left: 5px solid #00d2ff; margin-top: 15px; }
    section[data-testid="stSidebar"] { background-color: rgba(15, 23, 42, 0.8); border-right: 1px solid rgba(255, 255, 255, 0.1); }
    </style>
    """, unsafe_allow_html=True)

# --- 2. وظائف معالجة الملفات والتصدير ---
def get_word_download(text):
    doc = Document()
    doc.add_heading('AI Architect Pro - Analysis Report', 0)
    doc.add_paragraph(text)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def get_excel_download(text):
    try:
        from io import StringIO
        if "|" in text:
            lines = [l.strip() for l in text.split('\n') if "|" in l]
            if len(lines) > 2:
                df = pd.read_csv(StringIO('\n'.join(lines)), sep="|", skipinitialspace=True).dropna(axis=1, how='all')
                df.columns = [c.strip() for c in df.columns]
                out = io.BytesIO()
                with pd.ExcelWriter(out, engine='openpyxl') as writer:
                    df.to_excel(writer, index=False)
                return out.getvalue()
    except: return None
    return None

def process_office_file(file):
    ext = file.name.split('.')[-1].lower()
    content = ""
    try:
        if ext == 'docx':
            doc = Document(file)
            content = "\n".join([p.text for p in doc.paragraphs])
        elif ext == 'xlsx':
            df = pd.read_excel(file)
            content = f"Excel Data Summary of {file.name}:\n{df.to_string()}"
        elif ext == 'pptx':
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): content += shape.text + "\n"
    except Exception as e: content = f"Error reading {file.name}: {e}"
    return f"--- File Content: {file.name} ---\n{content}\n"

def encode_image(image):
    buffered = io.BytesIO()
    if image.mode in ("RGBA", "P"): image = image.convert("RGB")
    image.save(buffered, format="JPEG")
    return base64.b64encode(buffered.getvalue()).decode('utf-8')

# --- 3. محرك التوليد الذكي ---
def generate_response(provider, api_key, model_name, query, images=None):
    try:
        if provider == "Google Gemini":
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(model_name)
            res = model.generate_content([query] + (images if images else []))
            return res.text
        elif provider == "Groq (Ultra Fast)":
            client = Groq(api_key=api_key)
            if images and ("vision" in model_name.lower() or "3.2" in model_name.lower()):
                msgs = [{"role": "user", "content": [{"type": "text", "text": query}, {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encode_image(images[0])}"}}]}]
            else:
                msgs = [{"role": "user", "content": query}]
            res = client.chat.completions.create(messages=msgs, model=model_name)
            return res.choices[0].message.content
    except Exception as e:
        st.error(f"❌ Error: {str(e)}")
        return None

# --- 4. القائمة الجانبية (الاكتشاف التلقائي للموديلات) ---
with st.sidebar:
    st.markdown("<h2 style='text-align: center; color: #00d2ff;'>💎 Control Center</h2>", unsafe_allow_html=True)
    provider = st.selectbox("AI Provider:", ["Google Gemini", "Groq (Ultra Fast)"])
    api_key = st.text_input(f"{provider} API Key:", type="password")
    
    model_choice = None
    if api_key:
        with st.spinner("جاري جلب الموديلات المتاحة..."):
            try:
                if provider == "Google Gemini":
                    genai.configure(api_key=api_key)
                    models = [m.name.replace('models/', '') for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
                    model_choice = st.selectbox("اختر الموديل (Gemini):", models)
                else:
                    client = Groq(api_key=api_key)
                    groq_models = [m.id for m in client.models.list().data]
                    model_choice = st.selectbox("اختر الموديل (Groq):", groq_models)
                st.success(f"✅ متصل بـ {provider}")
            except Exception as e:
                st.error(f"❌ خطأ في الاتصال: تأكد من صحة الـ Key")

# --- 5. واجهة المستخدم الرئيسية والتابات ---
if api_key and model_choice:
    st.markdown("<h1 style='text-align: center;'>🪄 AI Architect <span style='color: #00d2ff;'>Ultimate</span></h1>", unsafe_allow_html=True)
    
    tabs = st.tabs(["📑 Doc Analyzer", "✨ Image Prompts", "📸 Vision Studio", "🧠 Universal Architect"])

    # --- Tab 1: Ultimate Doc Analyzer ---
    with tabs[0]:
        st.markdown("### 📑 PDF, Office, Code & Text Intelligence")
        docs = st.file_uploader("ارفع حتى 10 ملفات متنوعة", type=["pdf", "docx", "xlsx", "pptx", "txt", "py", "jpg", "png"], accept_multiple_files=True)
        
        payload_text = []
        payload_imgs = []
        if docs:
            for d in docs[:10]:
                ext = d.name.split('.')[-1].lower()
                if ext in ['docx', 'xlsx', 'pptx']: payload_text.append(process_office_file(d))
                elif ext in ['txt', 'py']: payload_text.append(f"--- File: {d.name} ---\n{d.getvalue().decode('utf-8')}\n")
                elif ext == 'pdf':
                    pdf_doc = fitz.open(stream=d.read(), filetype="pdf")
                    for page in pdf_doc:
                        pix = page.get_pixmap(matrix=fitz.Matrix(1,1))
                        payload_imgs.append(Image.open(io.BytesIO(pix.tobytes("png"))))
                elif ext in ['jpg', 'png', 'jpeg']: payload_imgs.append(Image.open(d))
            st.success(f"تم تحميل {len(docs[:10])} ملفات.")

        d_query = st.text_area("ما هي تعليماتك؟ (لخص، قارن، استخرج...)")
        if st.button("تحليل البيانات 🚀"):
            full_context = "".join(payload_text) + "\n" + d_query
            res = generate_response(provider, api_key, model_choice, full_context, payload_imgs if payload_imgs else None)
            if res:
                st.session_state['last_res'] = res
                st.markdown("### 🔍 النتيجة (جاهزة للنسخ):")
                st.code(res, language="markdown")
                
                col_d1, col_d2 = st.columns(2)
                col_d1.download_button("Download Word 📄", get_word_download(res), "Analysis_Report.docx")
                excel_data = get_excel_download(res)
                if excel_data:
                    col_d2.download_button("Download Excel 📊", excel_data, "Extracted_Data.xlsx")

    # --- Tab 2: Image Prompt Builder ---
    with tabs[1]:
        st.markdown("### ✍️ Image Prompts Builder")
        p_idea = st.text_area("صف فكرتك باللغة العربية:")
        p_target = st.selectbox("المنصة المستهدفة:", ["Midjourney", "DALL-E 3", "Leonardo AI"])
        if st.button("إنشاء برومبت احترافي ✨"):
            p_res = generate_response(provider, api_key, model_choice, f"Convert this idea into a high-detail English image prompt for {p_target}: {p_idea}")
            if p_res: st.code(p_res)

    # --- Tab 3: Vision Studio ---
    with tabs[2]:
        st.markdown("### 📸 Vision Intelligence")
        v_files = st.file_uploader("ارفع صوراً للتحليل", type=["jpg", "png", "jpeg"], accept_multiple_files=True, key="vision_up")
        v_q = st.text_input("ماذا تريد أن تعرف عن الصور؟")
        if st.button("تحليل الصور 👁️") and v_files:
            v_imgs = [Image.open(f) for f in v_files]
            v_res = generate_response(provider, api_key, model_choice, v_q if v_q else "Describe these images", v_imgs)
            if v_res: st.markdown(f'<div class="result-box">{v_res}</div>', unsafe_allow_html=True)

    # --- Tab 4: Universal Architect ---
    with tabs[3]:
        st.markdown("### 🧠 Universal Prompt Architect")
        u_idea = st.text_area("اكتب أي طلب وسأقوم ببناء 'أمر هندسي' متكامل له:")
        if st.button("Build Master Prompt 🔨"):
            u_res = generate_response(provider, api_key, model_choice, f"Create an expert level prompt including Role, Context and Task for: {u_idea}")
            if u_res: st.code(u_res)

else:
    st.markdown("<div style='text-align: center; padding: 100px;'><h2>👋 مرحباً بك في AI Architect</h2><p>يرجى إدخال الـ API Key في القائمة الجانبية للبدء</p></div>", unsafe_allow_html=True)
